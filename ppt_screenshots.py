"""Create 7-8 PPT slides with app screenshots embedded. Light background, professional."""
import sys,io
sys.stdout=io.TextIOWrapper(sys.stdout.buffer,encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches,Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

W=Inches(13.333);H=Inches(7.5)
prs=Presentation();prs.slide_width=W;prs.slide_height=H
BL=prs.slide_layouts[6]

# Colors - light professional theme
NAVY=RGBColor(0x0A,0x24,0x63); GOLD=RGBColor(0xB8,0x86,0x0B)
WHITE=RGBColor(0xFF,0xFF,0xFF); LGRAY=RGBColor(0xF5,0xF5,0xF5)
DARK=RGBColor(0x2C,0x2C,0x2C); MEDIUM=RGBColor(0x55,0x55,0x55)
ACCENT=RGBColor(0x1B,0x4F,0x8A)

def add_light_bg(sl):
    """Add light gradient background using a full-width shape."""
    from pptx.oxml.ns import qn
    bg = sl.background
    fill = bg.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = RGBColor(0xF8,0xFA,0xFC)
    fill.gradient_stops[0].position = 0
    fill.gradient_stops[1].color.rgb = RGBColor(0xE8,0xEE,0xF4)
    fill.gradient_stops[1].position = 1.0

def add_top_bar(sl, title_text):
    """Navy blue top bar with gold accent line."""
    bar = sl.shapes.add_shape(1, 0, 0, W, Inches(1.1))  # rectangle
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    # Gold accent line
    line = sl.shapes.add_shape(1, 0, Inches(1.1), W, Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()
    # Title text
    tx = sl.shapes.add_textbox(Inches(0.6), Inches(0.2), Inches(12), Inches(0.7))
    tf = tx.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run()
    r.text = title_text
    r.font.size = Pt(28); r.font.bold = True; r.font.name = 'Inter'; r.font.color.rgb = WHITE

def add_footer(sl, text="Bifacial PV Optimizer | NSUT, New Delhi | EndSem 2026"):
    tx = sl.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12), Inches(0.4))
    tf = tx.text_frame; tf.word_wrap = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = text; r.font.size = Pt(10); r.font.name = 'Inter'; r.font.color.rgb = MEDIUM

def txt(sl,t,l,tp,w,h,sz=16,c=DARK,b=False,al=PP_ALIGN.LEFT):
    tx=sl.shapes.add_textbox(Inches(l),Inches(tp),Inches(w),Inches(h))
    tf=tx.text_frame;tf.word_wrap=True;tf.paragraphs[0].alignment=al
    r=tf.paragraphs[0].add_run();r.text=t;r.font.size=Pt(sz);r.font.name='Inter';r.font.color.rgb=c;r.font.bold=b

def img(sl, path, l, tp, w, h):
    sl.shapes.add_picture(path, Inches(l), Inches(tp), Inches(w), Inches(h))

# ===================================================================
# SLIDE 1: PLATFORM OVERVIEW (Hero + Inputs)
# ===================================================================
sl = prs.slides.add_slide(BL); add_light_bg(sl)
add_top_bar(sl, 'Bifacial PV Optimizer \u2014 Platform Overview')
add_footer(sl)
img(sl, 'ppt_assets/ss1_hero.png', 0.5, 1.4, 7.5, 4.7)
txt(sl, 'Platform Features', 8.5, 1.5, 4.5, 0.5, sz=20, c=NAVY, b=True)
features = [
    '\u2022 Location-based solar data via\n  NASA POWER API',
    '\u2022 Configurable sweep parameters:\n  Height, Tilt, Step sizes',
    '\u2022 Interactive city search with\n  geocoding support',
    '\u2022 Date-specific hourly analysis\n  with IST timezone handling',
    '\u2022 Dark-themed professional UI\n  with real-time validation'
]
tx = sl.shapes.add_textbox(Inches(8.5), Inches(2.2), Inches(4.3), Inches(4.5))
tf = tx.text_frame; tf.word_wrap = True
for i, f in enumerate(features):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(10)
    r = p.add_run(); r.text = f; r.font.size = Pt(14); r.font.name = 'Inter'; r.font.color.rgb = DARK

# ===================================================================
# SLIDE 2: INPUT CONFIGURATION (Albedo + Buttons)
# ===================================================================
sl = prs.slides.add_slide(BL); add_light_bg(sl)
add_top_bar(sl, 'Input Configuration \u2014 Albedo & Surface Selection')
add_footer(sl)
img(sl, 'ppt_assets/ss2_albedo.png', 0.5, 1.4, 7.5, 4.7)
txt(sl, 'Surface Albedo Presets', 8.5, 1.5, 4.5, 0.5, sz=20, c=NAVY, b=True)
details = [
    '\u2022 5 preset surfaces:\n  Grass (0.20), Concrete (0.30),\n  Snow (0.82), Asphalt (0.12),\n  Custom (user-defined)',
    '\u2022 Custom albedo range: 0\u20131',
    '\u2022 "Fetch Irradiance" retrieves\n  hourly GHI from NASA POWER',
    '\u2022 "Run Global Analysis" triggers\n  full parametric sweep engine',
    '\u2022 Array azimuth configurable\n  (default 180\u00b0 for South-facing)'
]
tx = sl.shapes.add_textbox(Inches(8.5), Inches(2.2), Inches(4.3), Inches(4.5))
tf = tx.text_frame; tf.word_wrap = True
for i, d in enumerate(details):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(8)
    r = p.add_run(); r.text = d; r.font.size = Pt(14); r.font.name = 'Inter'; r.font.color.rgb = DARK

# ===================================================================
# SLIDE 3: NASA POWER DATA
# ===================================================================
sl = prs.slides.add_slide(BL); add_light_bg(sl)
add_top_bar(sl, 'NASA POWER Telemetry \u2014 Solar Irradiance Data')
add_footer(sl)
img(sl, 'ppt_assets/ss3_nasa.png', 0.5, 1.4, 7.5, 4.7)
txt(sl, 'Data Retrieved', 8.5, 1.5, 4.5, 0.5, sz=20, c=NAVY, b=True)
data_pts = [
    '\u2022 Source: NASA POWER API\n  (Hourly, 0.5\u00b0 resolution)',
    '\u2022 Parameters: ALLSKY_SFC_SW_DWN\n  (Global Horizontal Irradiance)',
    '\u2022 Erbs clearness-index model\n  decomposes GHI \u2192 DNI + DHI',
    '\u2022 Total GHI: ~5.43 kWh/m\u00b2/day\n  for Greater Noida',
    '\u2022 Hourly temperature data\n  for NOCT thermal derating'
]
tx = sl.shapes.add_textbox(Inches(8.5), Inches(2.2), Inches(4.3), Inches(4.5))
tf = tx.text_frame; tf.word_wrap = True
for i, d in enumerate(data_pts):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(8)
    r = p.add_run(); r.text = d; r.font.size = Pt(14); r.font.name = 'Inter'; r.font.color.rgb = DARK

# ===================================================================
# SLIDE 4: OPTIMIZATION RESULTS
# ===================================================================
sl = prs.slides.add_slide(BL); add_light_bg(sl)
add_top_bar(sl, 'Optimization Engine \u2014 Results Dashboard')
add_footer(sl)
img(sl, 'ppt_assets/ss4_results.png', 0.3, 1.4, 12.7, 5.3)
txt(sl, '\u2191 Optimal configuration identified: Height 450cm, Tilt 30\u00b0 \u2192 Max Energy; Height 450cm, Tilt 10\u00b0 \u2192 Max Rear Gain', 0.5, 6.7, 12, 0.4, sz=13, c=ACCENT, b=True, al=PP_ALIGN.CENTER)

# ===================================================================
# SLIDE 5: ENERGY COMPARISON CHART
# ===================================================================
sl = prs.slides.add_slide(BL); add_light_bg(sl)
add_top_bar(sl, 'Energy Yield Comparison \u2014 Top Configurations')
add_footer(sl)
img(sl, 'ppt_assets/ss5_energy.png', 0.3, 1.4, 12.7, 5.3)
txt(sl, '\u2191 Bar chart comparing total energy (kWh) across top-ranked configurations with optimal highlighted', 0.5, 6.7, 12, 0.4, sz=13, c=ACCENT, b=True, al=PP_ALIGN.CENTER)

# ===================================================================
# SLIDE 6: I-V AND P-V CURVES
# ===================================================================
sl = prs.slides.add_slide(BL); add_light_bg(sl)
add_top_bar(sl, 'Electrical Characteristics \u2014 I-V & P-V Curves')
add_footer(sl)
img(sl, 'ppt_assets/ss6_iv.png', 0.3, 1.3, 6.3, 4.2)
img(sl, 'ppt_assets/ss7_pv.png', 6.7, 1.3, 6.3, 4.2)
annotations = [
    '\u2022 I-V curves show current vs. voltage for top configurations with both front and rear irradiance',
    '\u2022 P-V curves identify Maximum Power Point (MPP) for each configuration',
    '\u2022 Higher effective irradiance \u2192 higher Isc and peak power; Voc remains relatively stable'
]
tx = sl.shapes.add_textbox(Inches(0.5), Inches(5.6), Inches(12), Inches(1.4))
tf = tx.text_frame; tf.word_wrap = True
for i, a in enumerate(annotations):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(3)
    r = p.add_run(); r.text = a; r.font.size = Pt(14); r.font.name = 'Inter'; r.font.color.rgb = DARK

# ===================================================================
# SLIDE 7: SCENARIO VARIATION
# ===================================================================
sl = prs.slides.add_slide(BL); add_light_bg(sl)
add_top_bar(sl, 'Scenario Variation \u2014 Individual Parameter Analysis')
add_footer(sl)
img(sl, 'ppt_assets/ss8_scenario.png', 0.3, 1.4, 12.7, 4.5)
annotations2 = [
    '\u2022 Vary Albedo: Rear gain ranges from <5% (asphalt) to >24% (aluminum) \u2014 most influential parameter',
    '\u2022 Vary Height: Logarithmic dependence; diminishing returns above 200-300cm',
    '\u2022 Vary Tilt: Max rear gain at low tilt (10\u00b0); front irradiance peaks near latitude tilt (28\u00b0)'
]
tx = sl.shapes.add_textbox(Inches(0.5), Inches(6), Inches(12), Inches(1.2))
tf = tx.text_frame; tf.word_wrap = True
for i, a in enumerate(annotations2):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(3)
    r = p.add_run(); r.text = a; r.font.size = Pt(14); r.font.name = 'Inter'; r.font.color.rgb = DARK

# ===================================================================
# SLIDE 8: FULL APP OVERVIEW
# ===================================================================
sl = prs.slides.add_slide(BL); add_light_bg(sl)
add_top_bar(sl, 'Complete Application \u2014 Full Stack Architecture')
add_footer(sl)
img(sl, 'ppt_assets/ss9_hero2.png', 0.3, 1.3, 6.3, 5)
txt(sl, 'Technology Stack', 7, 1.5, 5.5, 0.5, sz=22, c=NAVY, b=True)
stack_items = [
    ('\u25b6 Frontend:', 'HTML5, CSS3, JavaScript\n  Chart.js for visualization'),
    ('\u25b6 Backend:', 'Node.js + Express.js\n  RESTful API design'),
    ('\u25b6 Data Source:', 'NASA POWER API\n  Hourly GHI & Temperature'),
    ('\u25b6 Simulation:', 'MATLAB/Simulink engine\n  calculate_irradiance.m'),
    ('\u25b6 Analysis:', 'JS-based fallback engine\n  Parametric sweep + I-V/P-V'),
]
tx = sl.shapes.add_textbox(Inches(7), Inches(2.2), Inches(5.5), Inches(4.5))
tf = tx.text_frame; tf.word_wrap = True
for i, (label, desc) in enumerate(stack_items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(8)
    r1 = p.add_run(); r1.text = label; r1.font.size = Pt(16); r1.font.bold = True; r1.font.name = 'Inter'; r1.font.color.rgb = ACCENT
    r2 = p.add_run(); r2.text = ' ' + desc; r2.font.size = Pt(14); r2.font.name = 'Inter'; r2.font.color.rgb = DARK

prs.save('App_Screenshots_Slides.pptx')
print('=== SAVED: App_Screenshots_Slides.pptx (8 slides) ===')
