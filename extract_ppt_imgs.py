"""Extract images from midsem PPT for reuse."""
import sys, io, os
sys.stdout=io.TextIOWrapper(sys.stdout.buffer,encoding='utf-8')
from pptx import Presentation

os.makedirs('ppt_assets', exist_ok=True)
prs = Presentation('MidSem_PPT_BTP-2.pptx')

img_idx = 0
for si, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.shape_type == 13:  # PICTURE
            img = shape.image
            ext = img.content_type.split('/')[-1]
            if ext == 'jpeg': ext = 'jpg'
            fname = f'ppt_assets/slide{si+1}_img{img_idx}.{ext}'
            with open(fname, 'wb') as f:
                f.write(img.blob)
            print(f'Saved: {fname} ({shape.width/914400:.1f}x{shape.height/914400:.1f} in, pos={shape.left/914400:.1f},{shape.top/914400:.1f})')
            img_idx += 1

print(f'\nTotal images extracted: {img_idx}')
