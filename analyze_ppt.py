"""Analyze MidSem PPT structure for replication."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
import json

prs = Presentation('MidSem_PPT_BTP-2.pptx')

print(f"Slide width: {prs.slide_width}, height: {prs.slide_height}")
print(f"Slide width inches: {prs.slide_width/914400:.2f}, height inches: {prs.slide_height/914400:.2f}")
print(f"Total slides: {len(prs.slides)}")
print(f"Slide layouts available: {len(prs.slide_layouts)}")

for i, layout in enumerate(prs.slide_layouts):
    print(f"  Layout {i}: {layout.name}")

print("\n" + "="*80)
for idx, slide in enumerate(prs.slides):
    print(f"\n--- Slide {idx+1} ---")
    print(f"  Layout: {slide.slide_layout.name}")
    # Check background
    bg = slide.background
    fill = bg.fill
    print(f"  Background fill type: {fill.type}")
    
    for shape in slide.shapes:
        print(f"  Shape: {shape.shape_type}, name='{shape.name}', pos=({shape.left},{shape.top}), size=({shape.width},{shape.height})")
        if hasattr(shape, 'text') and shape.text:
            txt = shape.text[:120].replace('\n','\\n')
            print(f"    Text: '{txt}'")
        if shape.has_text_frame:
            for pi, para in enumerate(shape.text_frame.paragraphs):
                for ri, run in enumerate(para.runs):
                    font = run.font
                    color_str = ''
                    try:
                        if font.color and font.color.rgb:
                            color_str = str(font.color.rgb)
                    except:
                        pass
                    print(f"    P{pi}R{ri}: size={font.size}, bold={font.bold}, italic={font.italic}, color={color_str}, name={font.name}, text='{run.text[:60]}'")
        if shape.shape_type == 13:  # Picture
            print(f"    [IMAGE: {shape.image.content_type}]")
