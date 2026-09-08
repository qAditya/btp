"""EndSem PPT Part 1: Slides 1-10. Reuses midsem template backgrounds/logo."""
import sys,io
sys.stdout=io.TextIOWrapper(sys.stdout.buffer,encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

W=Inches(13.333); H=Inches(7.5)
prs=Presentation(); prs.slide_width=W; prs.slide_height=H
BL=prs.slide_layouts[6] # blank

NAVY=RGBColor(0x00,0x2B,0x5C); WHITE=RGBColor(0xFF,0xFF,0xFF)
GOLD=RGBColor(0xC4,0x9A,0x2A); DBLUE=RGBColor(0x1A,0x3C,0x6E)
BLACK=RGBColor(0,0,0); GRAY=RGBColor(0x33,0x33,0x33)

def add_bg(sl, img_path):
    sl.shapes.add_picture(img_path, 0, 0, W, H)

def add_logo(sl):
    sl.shapes.add_picture('ppt_assets/slide2_img1.png', Inches(12.4), Inches(0.04), Inches(0.88), Inches(0.88))

def add_title_bar(sl, text, top=Inches(0.1), sz=32, color=WHITE):
    tx=sl.shapes.add_textbox(Inches(0.5), top, Inches(12), Inches(1))
    tf=tx.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.LEFT
    r=p.add_run(); r.text=text; r.font.size=Pt(sz); r.font.bold=True; r.font.name='Times New Roman'; r.font.color.rgb=color

def add_text(sl, text, left, top, w, h, sz=18, color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    tx=sl.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf=tx.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=text; r.font.size=Pt(sz); r.font.name='Times New Roman'; r.font.color.rgb=color; r.font.bold=bold

def add_bullet_slide(sl, title, bullets, title_color=WHITE, body_color=WHITE, body_sz=19):
    add_title_bar(sl, title, color=title_color, sz=32)
    tx=sl.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(11.5), Inches(5.5))
    tf=tx.text_frame; tf.word_wrap=True
    for i,b in enumerate(bullets):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.space_after=Pt(8); p.alignment=PP_ALIGN.LEFT
        if isinstance(b, tuple):
            r=p.add_run(); r.text=b[0]; r.font.size=Pt(body_sz); r.font.bold=True; r.font.name='Times New Roman'; r.font.color.rgb=body_color
            r2=p.add_run(); r2.text=' '+b[1]; r2.font.size=Pt(body_sz); r2.font.name='Times New Roman'; r2.font.color.rgb=body_color
        else:
            r=p.add_run(); r.text='\u2022 '+b; r.font.size=Pt(body_sz); r.font.name='Times New Roman'; r.font.color.rgb=body_color

# ===== SLIDE 1: TITLE =====
sl=prs.slides.add_slide(BL)
add_bg(sl,'ppt_assets/slide1_img0.png')
add_text(sl,'OPTIMAL CONFIGURATION OF BIFACIAL\nPHOTOVOLTAIC MODULES USING\nPARAMETRIC ANALYSIS',1.5,1.5,10,2.5,sz=36,color=WHITE,bold=True,align=PP_ALIGN.CENTER)
add_text(sl,'End-Semester Presentation | B.Tech. Project',2.5,4.0,8,0.6,sz=20,color=GOLD,bold=True,align=PP_ALIGN.CENTER)
add_text(sl,'Ashish Upadhyay (2022UEE4521)  |  Pratyai Chakrabarty (2022UEE4586)\nAman (2022UEE4532)  |  Aditya Garg (2022UEE4503)',2,4.8,9,1,sz=16,color=WHITE,align=PP_ALIGN.CENTER)
add_text(sl,'Under the supervision of Dr. Astitva Kumar',3,5.8,7,0.5,sz=16,color=WHITE,bold=True,align=PP_ALIGN.CENTER)
add_text(sl,'Department of Electrical Engineering\nNetaji Subhas University of Technology, New Delhi',3,6.3,7,0.8,sz=14,color=WHITE,align=PP_ALIGN.CENTER)

# ===== SLIDE 2: OUTLINE =====
sl=prs.slides.add_slide(BL)
add_bg(sl,'ppt_assets/slide7_img6.jpg')
add_logo(sl)
add_title_bar(sl,'Presentation Outline',sz=34)
items=['1. Introduction & Motivation','2. Literature Survey & Gaps','3. Problem Statement & Objectives','4. Mathematical Formulation (Equations 1-11)','5. Simulink Model & Frontend Platform','6. Case Study: Greater Noida','7. Results & Analysis','8. Comparison with Literature','9. Conclusion & Key Findings','10. Future Scope','11. Publication & References']
tx=sl.shapes.add_textbox(Inches(1), Inches(1.4), Inches(5.5), Inches(5.5))
tf=tx.text_frame; tf.word_wrap=True
for i,item in enumerate(items):
    p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
    p.space_after=Pt(6)
    r=p.add_run(); r.text=item; r.font.size=Pt(22); r.font.name='Times New Roman'; r.font.color.rgb=WHITE; r.font.bold=True

# ===== SLIDE 3: INTRODUCTION =====
sl=prs.slides.add_slide(BL)
add_bg(sl,'ppt_assets/slide7_img6.jpg')
add_logo(sl)
add_title_bar(sl,'Introduction & Motivation',sz=34)
add_bullet_slide(sl,'',['Solar PV is the fastest-growing electricity source \u2014 global capacity exceeded 1,500 GW by end of 2024','India targets 500 GW non-fossil fuel capacity by 2030 under the National Solar Mission','Bifacial PV modules capture irradiance from both front and rear surfaces, achieving 5\u201330% higher energy yield','However, rear-side contribution is highly sensitive to tilt angle, ground albedo, and mounting height','Need for systematic, location-specific optimization tools to realize the full potential of bifacial technology','First industrial bifacial panels: 1984 \u2014 predicted to capture ~35% of global PV market by 2027'], body_sz=21)

# ===== SLIDE 4: LIT SURVEY =====
sl=prs.slides.add_slide(BL)
add_bg(sl,'ppt_assets/slide7_img6.jpg')
add_logo(sl)
add_title_bar(sl,'Literature Survey',sz=34)
bullets=[('Guerrero-Lemus et al. [7]:','Comprehensive bifacial PV technology review; need for standardized testing'),('Yusufoglu et al. [2]:','Annual performance analysis; up to 30% gain at 2m height; log-height dependence'),('Dincer & Ozer [3]:','Parametric analysis \u2014 albedo, tilt, height, mounting; 21.2% rear gain for aluminum'),('Sun et al. [10]:','Global optimization perspective; location-specific analysis essential'),('Ganesan et al. [4]:','n-type PERT bifacial under diverse albedo; 21.4% avg gain with aluminum'),('Pelaez et al. [5]:','Comparison of 5 bifacial irradiance models; ~20% bifacial gain validated')]
add_bullet_slide(sl,'',bullets,body_sz=20)

# ===== SLIDE 5: GAPS =====
sl=prs.slides.add_slide(BL)
add_bg(sl,'ppt_assets/slide7_img6.jpg')
add_logo(sl)
add_title_bar(sl,'Gaps Identified',sz=34)
bullets=[('Limited Combined Analysis:','Most studies evaluate tilt, albedo, height independently \u2014 not their combined interaction effect'),('Location-Specific Gap:','Many studies limited to European/American climates; India-specific validation lacking'),('Accessibility Barrier:','Current tools require specialized software (PVsyst, SAM) or commercial licenses'),('Single Optimization Objective:','Few studies consider dual objectives: max energy vs. max rear gain trade-off'),('No User-Friendly Platform:','No web-based tool exists for real-time location-specific bifacial PV optimization')]
add_bullet_slide(sl,'',bullets,body_sz=20)

# ===== SLIDE 6: PROBLEM STATEMENT =====
sl=prs.slides.add_slide(BL)
add_bg(sl,'ppt_assets/slide7_img6.jpg')
add_logo(sl)
add_title_bar(sl,'Problem Statement',sz=34)
tx=sl.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(2))
tf=tx.text_frame; tf.word_wrap=True
p=tf.paragraphs[0]; p.alignment=PP_ALIGN.CENTER
r=p.add_run(); r.text='"How can the rear-side energy contribution and overall performance of bifacial PV systems be maximized by optimally selecting tilt angle, ground albedo, and mounting height for a given geographical location?"'; r.font.size=Pt(24); r.font.bold=True; r.font.name='Times New Roman'; r.font.color.rgb=WHITE; r.font.italic=True
add_text(sl,'Research Objectives:',0.7,4.2,11,0.5,sz=22,color=GOLD,bold=True)
objs=['Develop comprehensive mathematical model for bifacial PV (front + rear irradiance)','Implement in MATLAB/Simulink for I-V / P-V characteristic simulation','Build user-oriented frontend platform with NASA POWER API integration','Perform parametric sweep for optimal configuration identification','Validate through case study and literature comparison']
tx2=sl.shapes.add_textbox(Inches(0.9), Inches(4.8), Inches(11), Inches(2.5))
tf2=tx2.text_frame; tf2.word_wrap=True
for i,o in enumerate(objs):
    p=tf2.paragraphs[0] if i==0 else tf2.add_paragraph()
    r=p.add_run(); r.text=f'{i+1}. {o}'; r.font.size=Pt(17); r.font.name='Times New Roman'; r.font.color.rgb=WHITE

# ===== SLIDE 7: MATH MODEL (Front) =====
sl=prs.slides.add_slide(BL)
add_bg(sl,'ppt_assets/slide7_img6.jpg')
add_logo(sl)
add_title_bar(sl,'Mathematical Formulation \u2014 Front Irradiance',sz=30)
eqs=[('Eq. 1:','G_tot = G_front + G_rear','Total effective irradiance on bifacial module'),('Eq. 2:','beam_tilted = DNI \u00d7 max(0, cos \u03b8\u1d62)','Direct beam on tilted surface'),('Eq. 3:','sky_vf = (1 + cos \u03b2) / 2','Sky view factor for tilted plane'),('Eq. 4:','gnd_vf = (1 \u2212 cos \u03b2) / 2','Ground view factor for tilted plane'),('Eq. 5:','G_front = beam + DHI\u00b7sky_vf + GHI\u00b7\u03c1\u00b7gnd_vf','Liu-Jordan isotropic model [1]')]
tx=sl.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(11.5), Inches(5.5))
tf=tx.text_frame; tf.word_wrap=True
for i,(label,eq,desc) in enumerate(eqs):
    p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
    p.space_after=Pt(14)
    r1=p.add_run(); r1.text=f'{label}  '; r1.font.size=Pt(18); r1.font.bold=True; r1.font.name='Times New Roman'; r1.font.color.rgb=GOLD
    r2=p.add_run(); r2.text=eq; r2.font.size=Pt(20); r2.font.name='Cambria Math'; r2.font.color.rgb=WHITE; r2.font.italic=True
    r3=p.add_run(); r3.text=f'\n       \u2192 {desc}'; r3.font.size=Pt(15); r3.font.name='Times New Roman'; r3.font.color.rgb=RGBColor(0xBB,0xBB,0xBB)

# ===== SLIDE 8: MATH MODEL (Rear) =====
sl=prs.slides.add_slide(BL)
add_bg(sl,'ppt_assets/slide7_img6.jpg')
add_logo(sl)
add_title_bar(sl,'Mathematical Formulation \u2014 Rear Irradiance & Shadow',sz=28)
eqs=[('Eq. 6:','\u03b3_profile = arctan(tan \u03b1_sun / cos(\u03b3_sun \u2212 \u03b3_panel))','Solar profile angle'),('Eq. 7:','shadow_lower = h / tan(\u03b3_profile)','Lower shadow edge on ground'),('Eq. 8:','shadow_upper = (h + W\u00b7sin\u03b2) / tan(\u03b3_profile)','Upper shadow edge'),('Eq. 9:','F_V = \u222b cos\u03b8\u2081\u00b7cos\u03b8\u2082 / (2r) dx','Shadow view factor [2]'),('Eq. 10:','G_r0 = \u03c1\u00b7DHI\u00b7rear_vf + \u03c1\u00b7(GHI\u2212DHI)\u00b7(rear_vf\u2212F_V)','Rear irradiance with shadow correction'),('Eq. 11:','G_rear = G_r0 \u00d7 \u03c6_bifaciality','Scaled by bifaciality factor (0.7)')]
tx=sl.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(11.5), Inches(5.5))
tf=tx.text_frame; tf.word_wrap=True
for i,(label,eq,desc) in enumerate(eqs):
    p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
    p.space_after=Pt(12)
    r1=p.add_run(); r1.text=f'{label}  '; r1.font.size=Pt(18); r1.font.bold=True; r1.font.name='Times New Roman'; r1.font.color.rgb=GOLD
    r2=p.add_run(); r2.text=eq; r2.font.size=Pt(19); r2.font.name='Cambria Math'; r2.font.color.rgb=WHITE; r2.font.italic=True
    r3=p.add_run(); r3.text=f'\n       \u2192 {desc}'; r3.font.size=Pt(14); r3.font.name='Times New Roman'; r3.font.color.rgb=RGBColor(0xBB,0xBB,0xBB)

# ===== SLIDE 9: SIMULINK =====
sl=prs.slides.add_slide(BL)
add_bg(sl,'ppt_assets/slide7_img6.jpg')
add_logo(sl)
add_title_bar(sl,'Simulink Model & Implementation',sz=32)
add_text(sl,'[Insert Simulink Model Screenshot Here]',1,2,5.5,3.5,sz=16,color=WHITE,align=PP_ALIGN.CENTER)
bullets=['calculate_irradiance MATLAB function block\nimplements Eqs. 1\u201311','PV Array block: single-diode model\n550 Wp bifacial module','V_ramp sweep: 0V to Voc for\ncomplete I-V / P-V curves','NOCT thermal model:\ncell temp derating','Simpson\'s rule numerical integration\nfor shadow view factor F_V (200 pts)']
tx=sl.shapes.add_textbox(Inches(7), Inches(1.5), Inches(5.5), Inches(5.5))
tf=tx.text_frame; tf.word_wrap=True
for i,b in enumerate(bullets):
    p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
    p.space_after=Pt(10)
    r=p.add_run(); r.text='\u2022 '+b; r.font.size=Pt(17); r.font.name='Times New Roman'; r.font.color.rgb=WHITE

# ===== SLIDE 10: PLATFORM =====
sl=prs.slides.add_slide(BL)
add_bg(sl,'ppt_assets/slide7_img6.jpg')
add_logo(sl)
add_title_bar(sl,'Frontend Platform Architecture',sz=32)
layers=[('User Interface','City selection, date range, sweep parameters, albedo surface type, optimization objective'),('Backend (Node.js/Express)','API routing, input validation, NASA POWER data retrieval, parametric sweep coordination'),('NASA POWER API','Hourly GHI, DHI, DNI, temperature data for any global location (0.5\u00b0 resolution)'),('Analysis Engine','Erbs decomposition, Liu-Jordan transposition, view factor computation, I-V/P-V generation'),('Results Dashboard','Optimal config tables, I-V/P-V charts, individual parameter variation graphs')]
tx=sl.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(11.5), Inches(5.5))
tf=tx.text_frame; tf.word_wrap=True
for i,(title,desc) in enumerate(layers):
    p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
    p.space_after=Pt(12)
    r1=p.add_run(); r1.text=f'\u25b6 {title}: '; r1.font.size=Pt(20); r1.font.bold=True; r1.font.name='Times New Roman'; r1.font.color.rgb=GOLD
    r2=p.add_run(); r2.text=desc; r2.font.size=Pt(18); r2.font.name='Times New Roman'; r2.font.color.rgb=WHITE

prs.save('_endsem_p1.pptx')
print('Part 1 saved (10 slides)')
