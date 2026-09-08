"""EndSem PPT Part 2 - fresh design. Slides 11-20."""
import sys,io
sys.stdout=io.TextIOWrapper(sys.stdout.buffer,encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches,Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

W=Inches(13.333);H=Inches(7.5)
prs=Presentation('_endsem_fresh_p1.pptx')
BL=prs.slide_layouts[6]
GOLD=RGBColor(0xC4,0x9A,0x2A);WHITE=RGBColor(255,255,255);LGRAY=RGBColor(0xBB,0xBB,0xBB)
BG_C='ppt_assets/bg_content.png';BG_TY='ppt_assets/bg_thankyou.png';BG_T='ppt_assets/bg_title.png'

def bg(sl,img=BG_C): sl.shapes.add_picture(img,0,0,W,H)
def ttl(sl,t,sz=32):
    tx=sl.shapes.add_textbox(Inches(0.6),Inches(0.15),Inches(12),Inches(0.9));tf=tx.text_frame;tf.word_wrap=True
    r=tf.paragraphs[0].add_run();r.text=t;r.font.size=Pt(sz);r.font.bold=True;r.font.name='Inter';r.font.color.rgb=WHITE
def txt(sl,t,l,tp,w,h,sz=18,c=WHITE,b=False,al=PP_ALIGN.LEFT):
    tx=sl.shapes.add_textbox(Inches(l),Inches(tp),Inches(w),Inches(h));tf=tx.text_frame;tf.word_wrap=True;tf.paragraphs[0].alignment=al
    r=tf.paragraphs[0].add_run();r.text=t;r.font.size=Pt(sz);r.font.name='Inter';r.font.color.rgb=c;r.font.bold=b
def tbl(sl,hdrs,rows,l,tp,w):
    t=sl.shapes.add_table(len(rows)+1,len(hdrs),Inches(l),Inches(tp),Inches(w),Inches(0.42*(len(rows)+1))).table
    for j,h in enumerate(hdrs):
        c=t.cell(0,j);c.text=h
        for p in c.text_frame.paragraphs:
            p.alignment=PP_ALIGN.CENTER
            for r in p.runs:r.font.size=Pt(12);r.font.bold=True;r.font.name='Inter'
    for i,row in enumerate(rows):
        for j,v in enumerate(row):
            c=t.cell(i+1,j);c.text=str(v)
            for p in c.text_frame.paragraphs:
                p.alignment=PP_ALIGN.CENTER
                for r in p.runs:r.font.size=Pt(11);r.font.name='Inter'

# ===== S11: CASE STUDY =====
sl=prs.slides.add_slide(BL);bg(sl)
ttl(sl,'Case Study: Greater Noida, India')
params=['\u2022 Location: 28.47\u00b0N, 77.50\u00b0E','\u2022 Date: 13 March 2026','\u2022 Height: 50\u2013450 cm (step 50, 9 levels)','\u2022 Tilt: 10\u00b0\u201350\u00b0 (step 10\u00b0, 5 levels)','\u2022 Albedo: 0.3 (Concrete surface)','\u2022 Azimuth: 180\u00b0 (South-facing)','\u2022 Bifaciality factor: 0.7','\u2022 Total GHI observed: 6.08 kWh/m\u00b2','\u2022 Peak GHI: 866.53 W/m\u00b2','\u2022 45 evaluation points (9 \u00d7 5)']
tx=sl.shapes.add_textbox(Inches(0.7),Inches(1.3),Inches(5.5),Inches(5.5));tf=tx.text_frame;tf.word_wrap=True
for i,p_text in enumerate(params):
    p=tf.paragraphs[0] if i==0 else tf.add_paragraph();p.space_after=Pt(6)
    r=p.add_run();r.text=p_text;r.font.size=Pt(19);r.font.name='Inter';r.font.color.rgb=WHITE
txt(sl,'[Insert Frontend Screenshot]',7,2,5.5,3.5,sz=15,c=LGRAY,al=PP_ALIGN.CENTER)

# ===== S12: MAX ENERGY =====
sl=prs.slides.add_slide(BL);bg(sl)
ttl(sl,'Results: Maximum Energy Configuration',sz=28)
txt(sl,'Optimization Objective: Maximize Total Energy Output',0.7,1.1,11,0.4,sz=17,c=GOLD,b=True)
tbl(sl,['Rank','Height (cm)','Tilt (\u00b0)','Surface','Energy (kWh)','Rear Gain (%)'],
    [['1','450','30','Concrete','3.16843','16.55'],['2','450','40','Concrete','3.16570','15.61'],['3','400','30','Concrete','3.16521','16.39'],['4','400','40','Concrete','3.16258','15.45'],['5','350','30','Concrete','3.16113','16.19']],0.7,1.7,11.5)
tx=sl.shapes.add_textbox(Inches(0.7),Inches(5.2),Inches(11),Inches(2));tf=tx.text_frame;tf.word_wrap=True
for i,ins in enumerate(['\u2714 Optimal: 30\u00b0 tilt, 450cm height \u2192 closely matches site latitude (28.47\u00b0N)','\u2714 Top 5 differ by < 0.25% in energy \u2192 robust performance near optimum','\u2714 16.55% rear gain demonstrates significant bifacial advantage over monofacial']):
    p=tf.paragraphs[0] if i==0 else tf.add_paragraph();p.space_after=Pt(4)
    r=p.add_run();r.text=ins;r.font.size=Pt(16);r.font.name='Inter';r.font.color.rgb=WHITE

# ===== S13: MAX REAR GAIN =====
sl=prs.slides.add_slide(BL);bg(sl)
ttl(sl,'Results: Maximum Rear Gain Configuration',sz=28)
txt(sl,'Optimization Objective: Maximize Rear-Side Contribution',0.7,1.1,11,0.4,sz=17,c=GOLD,b=True)
tbl(sl,['Rank','Height (cm)','Tilt (\u00b0)','Surface','Energy (kWh)','Rear Gain (%)'],
    [['1','450','10','Concrete','2.99422','17.90'],['2','400','10','Concrete','2.99130','17.76'],['3','350','10','Concrete','2.98755','17.57'],['4','300','10','Concrete','2.98258','17.33'],['5','450','20','Concrete','3.11109','17.33']],0.7,1.7,11.5)
tx=sl.shapes.add_textbox(Inches(0.7),Inches(5.2),Inches(11),Inches(2));tf=tx.text_frame;tf.word_wrap=True
for i,ins in enumerate(['\u2714 Max rear gain (17.90%) at 10\u00b0 tilt \u2014 low tilt maximizes ground view factor','\u2714 Trade-off: 5.5% lower total energy compared to max-energy configuration','\u2714 Dual optimization enables informed decision-making for system designers']):
    p=tf.paragraphs[0] if i==0 else tf.add_paragraph();p.space_after=Pt(4)
    r=p.add_run();r.text=ins;r.font.size=Pt(16);r.font.name='Inter';r.font.color.rgb=WHITE

# ===== S14: I-V / P-V =====
sl=prs.slides.add_slide(BL);bg(sl)
ttl(sl,'I-V and P-V Characteristics')
txt(sl,'[Insert I-V Curve Plot]',0.5,1.5,6,4,sz=15,c=LGRAY,al=PP_ALIGN.CENTER)
txt(sl,'[Insert P-V Curve Plot]',6.8,1.5,6,4,sz=15,c=LGRAY,al=PP_ALIGN.CENTER)
txt(sl,'\u2022 Higher effective irradiance \u2192 higher Isc and peak power on I-V curve\n\u2022 MPP shifts with installation configuration changes\n\u2022 Voc relatively stable (logarithmic irradiance dependence)',0.7,5.8,11,1.5,sz=16,c=WHITE)

# ===== S15: PARAMETER VARIATION =====
sl=prs.slides.add_slide(BL);bg(sl)
ttl(sl,'Individual Parameter Variation Analysis',sz=27)
for i,(t,body) in enumerate([('Albedo Effect','Fixed: tilt=20\u00b0, h=100cm\n\n\u2022 Aluminum (0.85): 24.72%\n\u2022 Concrete (0.30): ~12%\n\u2022 Asphalt (0.12): < 5%\n\u2022 Near-linear relationship\n\u2022 Most influential parameter\n\u2022 ~5\u00d7 variation across surfaces'),('Height Effect','Fixed: tilt=20\u00b0, albedo=0.18\n\n\u2022 Gain increases with height\n\u2022 6.61% at h=100cm\n\u2022 Logarithmic dependence\n\u2022 200-300cm sweet spot\n\u2022 Diminishing returns above\n\u2022 Consistent with [2]'),('Tilt Effect','Fixed: h=100cm, albedo=0.18\n\n\u2022 Max rear gain 7.24% at 10\u00b0\n\u2022 Decreases with higher tilt\n\u2022 Front vs rear trade-off\n\u2022 Combined optimization\n  needed for best results\n\u2022 Consistent with [3]')]):
    left=0.4+i*4.3
    txt(sl,t,left,1.2,4,0.5,sz=20,c=GOLD,b=True,al=PP_ALIGN.CENTER)
    txt(sl,body,left,1.8,4,5,sz=15,c=WHITE)

# ===== S16: COMPARISON =====
sl=prs.slides.add_slide(BL);bg(sl)
ttl(sl,'Validation: Comparison with Literature',sz=28)
tbl(sl,['Study','Scenario','Reported','This Study'],
    [['Dincer & Ozer [3]','Aluminum, tilt 20\u00b0','21.2%','24.72%'],
     ['Dincer & Ozer [3]','Height 40-100cm','4.1-4.5%','5.2-6.6%'],
     ['Yusufoglu et al. [2]','Albedo=0.2, 2m, Cairo','~13.8%','16.55%'],
     ['Ganesan et al. [4]','Aluminum surface','21.4%','24.72%'],
     ['Pelaez et al. [5]','High albedo config','~20%','17.9-24.7%']],0.7,1.3,11.5)
tx=sl.shapes.add_textbox(Inches(0.7),Inches(5),Inches(11),Inches(2));tf=tx.text_frame;tf.word_wrap=True
for i,ins in enumerate(['\u2714 Consistent trends and magnitudes validate the proposed framework','\u2714 Slight differences due to geographic location, irradiance conditions, modeling specifics','\u2714 Albedo sensitivity, height dependency, and tilt optimization all show good agreement']):
    p=tf.paragraphs[0] if i==0 else tf.add_paragraph();p.space_after=Pt(4)
    r=p.add_run();r.text=ins;r.font.size=Pt(17);r.font.name='Inter';r.font.color.rgb=WHITE

# ===== S17: CONCLUSION =====
sl=prs.slides.add_slide(BL);bg(sl)
ttl(sl,'Conclusion & Key Findings')
findings=['\u2714 Comprehensive bifacial PV optimization framework developed with validated math model (Eqs. 1\u201311)','\u2714 Optimal for Greater Noida: 30\u00b0 tilt, 450cm height \u2192 16.55% rear gain, 3.168 kWh','\u2714 Albedo is dominant parameter: ~5\u00d7 variation in rear gain across surface types','\u2714 Height shows logarithmic benefit; 200\u2013300cm is the practical sweet spot','\u2714 Dual optimization reveals quantifiable trade-off: max energy vs max rear gain','\u2714 Web platform makes bifacial optimization accessible without commercial software','\u2714 Results validated against Dincer & Ozer, Yusufoglu et al., Ganesan et al., Pelaez et al.']
tx=sl.shapes.add_textbox(Inches(0.7),Inches(1.3),Inches(11.5),Inches(5.5));tf=tx.text_frame;tf.word_wrap=True
for i,f in enumerate(findings):
    p=tf.paragraphs[0] if i==0 else tf.add_paragraph();p.space_after=Pt(10)
    r=p.add_run();r.text=f;r.font.size=Pt(20);r.font.name='Inter';r.font.color.rgb=WHITE

# ===== S18: FUTURE WORK =====
sl=prs.slides.add_slide(BL);bg(sl)
ttl(sl,'Future Scope')
items=['Annual energy yield simulation (365-day aggregated analysis)','Multi-row array modeling with inter-row shading effects','Machine learning-based rapid optimization (surrogate models)','Techno-economic analysis with LCOE integration','Real-time IoT monitoring & adaptive optimization','GIS integration for regional deployment planning','Single/dual-axis tracking system comparison','Anisotropic sky models (Perez / Hay-Davies)']
tx=sl.shapes.add_textbox(Inches(0.7),Inches(1.3),Inches(11.5),Inches(5.5));tf=tx.text_frame;tf.word_wrap=True
for i,item in enumerate(items):
    p=tf.paragraphs[0] if i==0 else tf.add_paragraph();p.space_after=Pt(9)
    r=p.add_run();r.text=f'  {i+1}.  {item}';r.font.size=Pt(21);r.font.name='Inter';r.font.color.rgb=WHITE;r.font.bold=True

# ===== S19: REFERENCES =====
sl=prs.slides.add_slide(BL);bg(sl)
ttl(sl,'References')
refs_l=['[1] B.Y.H. Liu, R.C. Jordan, Solar Energy, vol. 4, 1960','[2] U.A. Yusufoglu et al., IEEE J. Photovolt., vol. 5, 2015','[3] F. Dincer, E. Ozer, Energies, vol. 18, 2025','[4] K. Ganesan et al., Solar Energy, vol. 252, 2023','[5] S.A. Pelaez et al., IEEE J. Photovolt., vol. 9, 2019','[6] N. Riedel-Lyngskjaer et al., Solar Energy, vol. 231, 2022','[7] R. Guerrero-Lemus et al., RSER, vol. 60, 2016','[8] M. Alam et al., SETA, vol. 57, 2023','[9] C. Deline et al., NREL Tech. Rep., 2019','[10] X. Sun et al., Appl. Energy, vol. 212, 2018','[11] A. Basak et al., Appl. Energy, vol. 345, 2025']
refs_r=['[12] U. Peter, M. Novak, J. Sustain. Energy Syst., 2025','[13] M.H. Aksoy, H.A. Ceylan, RERJ, vol. 13, 2023','[14] N. Baghel et al., Solar Energy, vol. 250, 2023','[15] I.N. Atalay et al., Solar Energy, vol. 246, 2023','[16] D.S. Braga et al., REES, vol. 8, 2023','[17] A.F. Almarshoud et al., Energies, vol. 17, 2024','[18] M.T. Patel et al., IEEE J. Photovolt., vol. 11, 2021','[19] J.A. Duffie, W.A. Beckman, Wiley, 2013','[20] B. Marion et al., IEEE PVSC, 2017','[21] P.K. Sahu et al., Solar Energy, vol. 262, 2023','[22] C. Ghenai et al., Solar Energy, vol. 223, 2021']
for refs,left in [(refs_l,0.5),(refs_r,6.8)]:
    tx=sl.shapes.add_textbox(Inches(left),Inches(1.2),Inches(6),Inches(5.8));tf=tx.text_frame;tf.word_wrap=True
    for i,ref in enumerate(refs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph();p.space_after=Pt(3)
        r=p.add_run();r.text=ref;r.font.size=Pt(13);r.font.name='Inter';r.font.color.rgb=WHITE

# ===== S20: THANK YOU =====
sl=prs.slides.add_slide(BL);bg(sl,BG_TY)
txt(sl,'Thank You!',2,2,9,1.5,sz=52,c=WHITE,b=True,al=PP_ALIGN.CENTER)
txt(sl,'Questions & Discussion',3,3.8,7,0.8,sz=28,c=GOLD,b=True,al=PP_ALIGN.CENTER)
txt(sl,'Publication:  "Performance Analysis of Bifacial PV System Based on Different Albedos"\nAccepted at IC2PCT 2026 (International Conference on Clean and Prospective Computing Technologies)',1.5,5,10,0.8,sz=15,c=WHITE,al=PP_ALIGN.CENTER)
txt(sl,'Supervisor: Dr. Astitva Kumar  |  Dept. of Electrical Engineering, NSUT, New Delhi',2,6.2,9,0.4,sz=13,c=LGRAY,al=PP_ALIGN.CENTER)

prs.save('EndSem_PPT_BTP_Final.pptx')
print('=== SAVED: EndSem_PPT_BTP_Final.pptx (20 slides) ===')
