"""EndSem PPT Part 2: Slides 11-20. Appends to _endsem_p1.pptx."""
import sys,io
sys.stdout=io.TextIOWrapper(sys.stdout.buffer,encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
# table alignment handled inline

W=Inches(13.333); H=Inches(7.5)
prs=Presentation('_endsem_p1.pptx')
BL=prs.slide_layouts[6]
GOLD=RGBColor(0xC4,0x9A,0x2A); WHITE=RGBColor(0xFF,0xFF,0xFF); BLACK=RGBColor(0,0,0)
LGRAY=RGBColor(0xBB,0xBB,0xBB)

def add_bg(sl): sl.shapes.add_picture('ppt_assets/slide7_img6.jpg',0,0,W,H)
def add_logo(sl): sl.shapes.add_picture('ppt_assets/slide2_img1.png',Inches(12.4),Inches(0.04),Inches(0.88),Inches(0.88))
def title(sl,t,sz=32):
    tx=sl.shapes.add_textbox(Inches(0.5),Inches(0.1),Inches(12),Inches(1));tf=tx.text_frame;tf.word_wrap=True
    r=tf.paragraphs[0].add_run();r.text=t;r.font.size=Pt(sz);r.font.bold=True;r.font.name='Times New Roman';r.font.color.rgb=WHITE
def txt(sl,t,l,tp,w,h,sz=18,c=WHITE,b=False,al=PP_ALIGN.LEFT):
    tx=sl.shapes.add_textbox(Inches(l),Inches(tp),Inches(w),Inches(h));tf=tx.text_frame;tf.word_wrap=True
    r=tf.paragraphs[0].add_run();r.text=t;r.font.size=Pt(sz);r.font.name='Times New Roman';r.font.color.rgb=c;r.font.bold=b;tf.paragraphs[0].alignment=al

def add_table(sl, headers, rows, left, top, width, row_h):
    tbl=sl.shapes.add_table(len(rows)+1,len(headers),Inches(left),Inches(top),Inches(width),Inches(row_h*(len(rows)+1))).table
    for j,h in enumerate(headers):
        c=tbl.cell(0,j);c.text=h
        for p in c.text_frame.paragraphs:
            p.alignment=PP_ALIGN.CENTER
            for r in p.runs: r.font.size=Pt(12);r.font.bold=True;r.font.name='Times New Roman'
    for i,row in enumerate(rows):
        for j,v in enumerate(row):
            c=tbl.cell(i+1,j);c.text=str(v)
            for p in c.text_frame.paragraphs:
                p.alignment=PP_ALIGN.CENTER
                for r in p.runs: r.font.size=Pt(11);r.font.name='Times New Roman'
    return tbl

# ===== SLIDE 11: CASE STUDY =====
sl=prs.slides.add_slide(BL);add_bg(sl);add_logo(sl)
title(sl,'Case Study: Greater Noida, India')
params=['\u2022 Location: 28.47\u00b0N, 77.50\u00b0E','\u2022 Date: 13 March 2026','\u2022 Height sweep: 50\u2013450 cm (step 50)','\u2022 Tilt sweep: 10\u00b0\u201350\u00b0 (step 10\u00b0)','\u2022 Albedo: 0.3 (Concrete)','\u2022 Azimuth: 180\u00b0 (South-facing)','\u2022 Bifaciality: 0.7','\u2022 Total GHI: 6.08 kWh/m\u00b2','\u2022 Peak GHI: 866.53 W/m\u00b2','\u2022 45 evaluation points (9\u00d75)']
tx=sl.shapes.add_textbox(Inches(0.7),Inches(1.4),Inches(6),Inches(5.5));tf=tx.text_frame;tf.word_wrap=True
for i,p_text in enumerate(params):
    p=tf.paragraphs[0] if i==0 else tf.add_paragraph();p.space_after=Pt(6)
    r=p.add_run();r.text=p_text;r.font.size=Pt(19);r.font.name='Times New Roman';r.font.color.rgb=WHITE
txt(sl,'[Insert Frontend Screenshot]',7.5,2,5,3,sz=16,c=LGRAY,al=PP_ALIGN.CENTER)

# ===== SLIDE 12: RESULTS TABLE - MAX ENERGY =====
sl=prs.slides.add_slide(BL);add_bg(sl);add_logo(sl)
title(sl,'Results: Maximum Energy Configuration',sz=30)
txt(sl,'Optimization Objective: Maximize Total Energy Output',0.7,1.2,11,0.5,sz=18,c=GOLD,b=True)
add_table(sl,['Rank','Height (cm)','Tilt (\u00b0)','Surface','Energy (kWh)','Rear Gain (%)'],
    [['1','450','30','Concrete','3.16843','16.55'],['2','450','40','Concrete','3.16570','15.61'],['3','400','30','Concrete','3.16521','16.39'],['4','400','40','Concrete','3.16258','15.45'],['5','350','30','Concrete','3.16113','16.19']],
    0.7,1.8,11.5,0.45)
tx=sl.shapes.add_textbox(Inches(0.7),Inches(5),Inches(11),Inches(2));tf=tx.text_frame;tf.word_wrap=True
insights=['\u2713 Optimal: 30\u00b0 tilt, 450cm height \u2192 matches latitude (28.47\u00b0N)','\u2713 Top 5 configs differ by < 0.25% in energy \u2192 robust near optimum','\u2713 Rear gain of 16.55% demonstrates significant bifacial advantage']
for i,ins in enumerate(insights):
    p=tf.paragraphs[0] if i==0 else tf.add_paragraph();p.space_after=Pt(4)
    r=p.add_run();r.text=ins;r.font.size=Pt(17);r.font.name='Times New Roman';r.font.color.rgb=WHITE

# ===== SLIDE 13: RESULTS TABLE - MAX REAR GAIN =====
sl=prs.slides.add_slide(BL);add_bg(sl);add_logo(sl)
title(sl,'Results: Maximum Rear Gain Configuration',sz=30)
txt(sl,'Optimization Objective: Maximize Rear-Side Contribution',0.7,1.2,11,0.5,sz=18,c=GOLD,b=True)
add_table(sl,['Rank','Height (cm)','Tilt (\u00b0)','Surface','Energy (kWh)','Rear Gain (%)'],
    [['1','450','10','Concrete','2.99422','17.90'],['2','400','10','Concrete','2.99130','17.76'],['3','350','10','Concrete','2.98755','17.57'],['4','300','10','Concrete','2.98258','17.33'],['5','450','20','Concrete','3.11109','17.33']],
    0.7,1.8,11.5,0.45)
tx=sl.shapes.add_textbox(Inches(0.7),Inches(5),Inches(11),Inches(2));tf=tx.text_frame;tf.word_wrap=True
insights=['\u2713 Max rear gain (17.90%) at 10\u00b0 tilt \u2014 low tilt \u2192 larger ground view factor','\u2713 Energy trade-off: 5.5% lower total energy vs max-energy config','\u2713 Dual optimization enables informed design decisions']
for i,ins in enumerate(insights):
    p=tf.paragraphs[0] if i==0 else tf.add_paragraph();p.space_after=Pt(4)
    r=p.add_run();r.text=ins;r.font.size=Pt(17);r.font.name='Times New Roman';r.font.color.rgb=WHITE

# ===== SLIDE 14: I-V / P-V =====
sl=prs.slides.add_slide(BL);add_bg(sl);add_logo(sl)
title(sl,'I-V and P-V Characteristics')
txt(sl,'[Insert I-V Curve Plot]',0.5,1.5,6,4,sz=16,c=LGRAY,al=PP_ALIGN.CENTER)
txt(sl,'[Insert P-V Curve Plot]',6.8,1.5,6,4,sz=16,c=LGRAY,al=PP_ALIGN.CENTER)
txt(sl,'\u2022 Higher effective irradiance \u2192 higher Isc and peak power\n\u2022 MPP shifts with installation configuration\n\u2022 Voc relatively stable across configs (logarithmic irradiance dependence)',0.7,5.8,11,1.5,sz=16,c=WHITE)

# ===== SLIDE 15: PARAMETER VARIATION =====
sl=prs.slides.add_slide(BL);add_bg(sl);add_logo(sl)
title(sl,'Individual Parameter Variation Analysis',sz=28)
cols=[('Albedo Effect','Fixed: tilt=20\u00b0, h=100cm\n\n\u2022 Aluminum (0.85): 24.72% rear gain\n\u2022 Concrete (0.30): ~12%\n\u2022 Asphalt (0.12): < 5%\n\u2022 Near-linear relationship\n\u2022 Most influential parameter'),
     ('Height Effect','Fixed: tilt=20\u00b0, albedo=0.18\n\n\u2022 Rear gain increases with height\n\u2022 6.61% at h=100cm\n\u2022 Diminishing returns (log)\n\u2022 200-300cm practical sweet spot'),
     ('Tilt Effect','Fixed: h=100cm, albedo=0.18\n\n\u2022 Max rear gain 7.24% at 10\u00b0\n\u2022 Decreases with higher tilt\n\u2022 Trade-off: front vs rear\n\u2022 Combined optimization needed')]
for i,(t,body) in enumerate(cols):
    left=0.5+i*4.2
    txt(sl,t,left,1.3,3.8,0.5,sz=20,c=GOLD,b=True,al=PP_ALIGN.CENTER)
    txt(sl,body,left,1.9,3.8,5,sz=15,c=WHITE)

# ===== SLIDE 16: COMPARISON =====
sl=prs.slides.add_slide(BL);add_bg(sl);add_logo(sl)
title(sl,'Comparison with Published Literature',sz=30)
add_table(sl,['Study','Scenario','Reported','This Study'],
    [['Dincer & Ozer [3]','Aluminum, tilt 20\u00b0','21.2%','24.72%'],
     ['Dincer & Ozer [3]','Height 40-100cm','4.1-4.5%','5.2-6.6%'],
     ['Yusufoglu et al. [2]','Albedo=0.2, 2m, Cairo','~13.8%','16.55%'],
     ['Ganesan et al. [4]','Aluminum surface','21.4%','24.72%'],
     ['Pelaez et al. [5]','High albedo config','~20%','17.9-24.7%']],
    0.7,1.5,11.5,0.45)
txt(sl,'\u2713 Consistent trends and magnitudes validate the proposed framework\n\u2713 Slight differences due to geographic location, irradiance conditions, and modeling methodology\n\u2713 Albedo sensitivity, height dependency, and tilt optimization all show good agreement',0.7,5.2,11,2,sz=17,c=WHITE)

# ===== SLIDE 17: CONCLUSION =====
sl=prs.slides.add_slide(BL);add_bg(sl);add_logo(sl)
title(sl,'Conclusion & Key Findings',sz=32)
findings=['\u2713 Developed comprehensive bifacial PV optimization framework with validated math model (Eqs. 1-11)','\u2713 Optimal config for Greater Noida: 30\u00b0 tilt, 450cm height \u2192 16.55% rear gain, 3.168 kWh','\u2713 Albedo is the dominant parameter: 5\u00d7 variation in rear gain across surface types','\u2713 Height shows logarithmic benefit; 200-300cm is the practical sweet spot','\u2713 Dual optimization (max energy vs max rear gain) reveals clear, quantifiable trade-offs','\u2713 Web-based platform makes bifacial optimization accessible without specialized software','\u2713 Results validated against Dincer & Ozer, Yusufoglu et al., Ganesan et al., Pelaez et al.']
tx=sl.shapes.add_textbox(Inches(0.7),Inches(1.4),Inches(11.5),Inches(5.5));tf=tx.text_frame;tf.word_wrap=True
for i,f in enumerate(findings):
    p=tf.paragraphs[0] if i==0 else tf.add_paragraph();p.space_after=Pt(10)
    r=p.add_run();r.text=f;r.font.size=Pt(20);r.font.name='Times New Roman';r.font.color.rgb=WHITE

# ===== SLIDE 18: FUTURE WORK =====
sl=prs.slides.add_slide(BL)
sl.shapes.add_picture('ppt_assets/slide17_img28.png',0,0,W,H)
add_logo(sl)
title(sl,'Future Scope')
items=['Annual energy yield simulation (365-day analysis)','Multi-row array modeling with inter-row shading','Machine learning-based rapid optimization','Techno-economic analysis (LCOE integration)','Real-time IoT monitoring & adaptive optimization','GIS integration for regional deployment planning','Single/dual-axis tracking system comparison','Anisotropic sky models (Perez/Hay-Davies)']
tx=sl.shapes.add_textbox(Inches(0.7),Inches(1.4),Inches(11.5),Inches(5.5));tf=tx.text_frame;tf.word_wrap=True
for i,item in enumerate(items):
    p=tf.paragraphs[0] if i==0 else tf.add_paragraph();p.space_after=Pt(8)
    r=p.add_run();r.text=f'{i+1}. {item}';r.font.size=Pt(21);r.font.name='Times New Roman';r.font.color.rgb=WHITE;r.font.bold=True

# ===== SLIDE 19: REFERENCES =====
sl=prs.slides.add_slide(BL);add_bg(sl);add_logo(sl)
title(sl,'References',sz=32)
refs=['[1] B.Y.H. Liu, R.C. Jordan, Solar Energy, vol. 4, 1960','[2] U.A. Yusufoglu et al., IEEE J. Photovolt., vol. 5, 2015','[3] F. Dincer, E. Ozer, Energies, vol. 18, 2025','[4] K. Ganesan et al., Solar Energy, vol. 252, 2023','[5] S.A. Pelaez et al., IEEE J. Photovolt., vol. 9, 2019','[6] N. Riedel-Lyngskjaer et al., Solar Energy, vol. 231, 2022','[7] R. Guerrero-Lemus et al., RSER, vol. 60, 2016','[8] M. Alam et al., SETA, vol. 57, 2023','[9] C. Deline et al., NREL Tech. Rep., 2019','[10] X. Sun et al., Appl. Energy, vol. 212, 2018','[11] A. Basak et al., Appl. Energy, vol. 345, 2025']
tx=sl.shapes.add_textbox(Inches(0.5),Inches(1.3),Inches(6),Inches(5.5));tf=tx.text_frame;tf.word_wrap=True
for i,ref in enumerate(refs):
    p=tf.paragraphs[0] if i==0 else tf.add_paragraph();p.space_after=Pt(3)
    r=p.add_run();r.text=ref;r.font.size=Pt(14);r.font.name='Times New Roman';r.font.color.rgb=WHITE
refs2=['[12] U. Peter, M. Novak, J. Sustain. Energy Syst., 2025','[13] M.H. Aksoy, H.A. Ceylan, RERJ, vol. 13, 2023','[14] N. Baghel et al., Solar Energy, vol. 250, 2023','[15] I.N. Atalay et al., Solar Energy, vol. 246, 2023','[16] D.S. Braga et al., REES, vol. 8, 2023','[17] A.F. Almarshoud et al., Energies, vol. 17, 2024','[18] M.T. Patel et al., IEEE J. Photovolt., vol. 11, 2021','[19] J.A. Duffie, W.A. Beckman, Wiley, 2013','[20] B. Marion et al., IEEE PVSC, 2017','[21] P.K. Sahu et al., Solar Energy, vol. 262, 2023','[22] C. Ghenai et al., Solar Energy, vol. 223, 2021']
tx2=sl.shapes.add_textbox(Inches(6.8),Inches(1.3),Inches(6),Inches(5.5));tf2=tx2.text_frame;tf2.word_wrap=True
for i,ref in enumerate(refs2):
    p=tf2.paragraphs[0] if i==0 else tf2.add_paragraph();p.space_after=Pt(3)
    r=p.add_run();r.text=ref;r.font.size=Pt(14);r.font.name='Times New Roman';r.font.color.rgb=WHITE

# ===== SLIDE 20: THANK YOU =====
sl=prs.slides.add_slide(BL)
sl.shapes.add_picture('ppt_assets/slide1_img0.png',0,0,W,H)
sl.shapes.add_picture('ppt_assets/slide20_img31.png',Inches(6),Inches(0.3),Inches(1.4),Inches(1.4))
txt(sl,'Thank You!',2,2.5,9,1.5,sz=48,c=WHITE,b=True,al=PP_ALIGN.CENTER)
txt(sl,'Questions & Discussion',3,4,7,0.8,sz=28,c=GOLD,b=True,al=PP_ALIGN.CENTER)
txt(sl,'Publication: "Performance Analysis of Bifacial PV System Based on Different Albedos"\nAccepted at IC2PCT 2026',2,5.2,9,1,sz=16,c=WHITE,al=PP_ALIGN.CENTER)
txt(sl,'Dr. Astitva Kumar (Supervisor) | Dept. of Electrical Engineering, NSUT',2.5,6.3,8,0.5,sz=14,c=LGRAY,al=PP_ALIGN.CENTER)

prs.save('EndSem_PPT_BTP.pptx')
print('=== FINAL PPT SAVED: EndSem_PPT_BTP.pptx (20 slides) ===')
