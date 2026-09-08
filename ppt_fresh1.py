"""EndSem PPT - fully original design with custom AI-generated backgrounds."""
import sys,io
sys.stdout=io.TextIOWrapper(sys.stdout.buffer,encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

W=Inches(13.333);H=Inches(7.5)
prs=Presentation();prs.slide_width=W;prs.slide_height=H
BL=prs.slide_layouts[6]
GOLD=RGBColor(0xC4,0x9A,0x2A);WHITE=RGBColor(255,255,255);LGRAY=RGBColor(0xBB,0xBB,0xBB)
BG_T='ppt_assets/bg_title.png';BG_C='ppt_assets/bg_content.png';BG_TY='ppt_assets/bg_thankyou.png'

def bg(sl,img=BG_C): sl.shapes.add_picture(img,0,0,W,H)
def ttl(sl,t,sz=32,top=0.15,c=WHITE):
    tx=sl.shapes.add_textbox(Inches(0.6),Inches(top),Inches(12),Inches(0.9));tf=tx.text_frame;tf.word_wrap=True
    r=tf.paragraphs[0].add_run();r.text=t;r.font.size=Pt(sz);r.font.bold=True;r.font.name='Inter';r.font.color.rgb=c
def txt(sl,t,l,tp,w,h,sz=18,c=WHITE,b=False,al=PP_ALIGN.LEFT,fn='Inter'):
    tx=sl.shapes.add_textbox(Inches(l),Inches(tp),Inches(w),Inches(h));tf=tx.text_frame;tf.word_wrap=True;tf.paragraphs[0].alignment=al
    r=tf.paragraphs[0].add_run();r.text=t;r.font.size=Pt(sz);r.font.name=fn;r.font.color.rgb=c;r.font.bold=b
def bullets(sl,items,l=0.7,tp=1.5,w=11.5,h=5.5,sz=19,c=WHITE):
    tx=sl.shapes.add_textbox(Inches(l),Inches(tp),Inches(w),Inches(h));tf=tx.text_frame;tf.word_wrap=True
    for i,item in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph();p.space_after=Pt(8);p.alignment=PP_ALIGN.LEFT
        if isinstance(item,tuple):
            r=p.add_run();r.text=item[0];r.font.size=Pt(sz);r.font.bold=True;r.font.name='Inter';r.font.color.rgb=GOLD
            r2=p.add_run();r2.text=' '+item[1];r2.font.size=Pt(sz);r2.font.name='Inter';r2.font.color.rgb=c
        else:
            r=p.add_run();r.text='\u2022 '+item;r.font.size=Pt(sz);r.font.name='Inter';r.font.color.rgb=c
def tbl(sl,hdrs,rows,l,tp,w):
    t=sl.shapes.add_table(len(rows)+1,len(hdrs),Inches(l),Inches(tp),Inches(w),Inches(0.4*(len(rows)+1))).table
    for j,h in enumerate(hdrs):
        c=t.cell(0,j);c.text=h
        for p in c.text_frame.paragraphs:
            p.alignment=PP_ALIGN.CENTER
            for r in p.runs:r.font.size=Pt(12);r.font.bold=True;r.font.name='Inter'
    for i,row in enumerate(rows):
        for j,v in enumerate(row):
            c=t.cell(i+1,j);c.text=str(v)
            for p in c.text_frame.paragraphs:
                p.alignment=PP_ALIGN.CENTER
                for r in p.runs:r.font.size=Pt(11);r.font.name='Inter'

# ===== S1: TITLE =====
sl=prs.slides.add_slide(BL);bg(sl,BG_T)
txt(sl,'END-SEMESTER PRESENTATION',2,0.8,9,0.5,sz=16,c=GOLD,b=True,al=PP_ALIGN.CENTER)
txt(sl,'Optimal Configuration of Bifacial\nPhotovoltaic Modules Using\nParametric Analysis',1.5,1.5,10,2.5,sz=40,c=WHITE,b=True,al=PP_ALIGN.CENTER,fn='Inter')
txt(sl,'B.Tech. Project | Department of Electrical Engineering',2.5,4.2,8,0.5,sz=16,c=LGRAY,al=PP_ALIGN.CENTER)
txt(sl,'Ashish Upadhyay (2022UEE4521)   \u2022   Pratyai Chakrabarty (2022UEE4586)\nAman (2022UEE4532)   \u2022   Aditya Garg (2022UEE4503)',2,5.0,9,0.8,sz=15,c=WHITE,al=PP_ALIGN.CENTER)
txt(sl,'Supervisor: Dr. Astitva Kumar, Assistant Professor',3,5.9,7,0.4,sz=14,c=GOLD,b=True,al=PP_ALIGN.CENTER)
txt(sl,'Netaji Subhas University of Technology (NSUT), New Delhi  |  May 2026',2.5,6.5,8,0.4,sz=13,c=LGRAY,al=PP_ALIGN.CENTER)

# ===== S2: OUTLINE =====
sl=prs.slides.add_slide(BL);bg(sl)
ttl(sl,'Presentation Outline',sz=34)
items=['Introduction & Motivation','Literature Survey','Research Gaps','Problem Statement & Objectives','Mathematical Formulation','Simulink Model & Platform Architecture','Case Study: Greater Noida','Results & Parametric Analysis','Validation with Literature','Conclusion & Future Scope','Publication & References']
tx=sl.shapes.add_textbox(Inches(1),Inches(1.3),Inches(6),Inches(5.5));tf=tx.text_frame;tf.word_wrap=True
for i,item in enumerate(items):
    p=tf.paragraphs[0] if i==0 else tf.add_paragraph();p.space_after=Pt(7)
    r=p.add_run();r.text=f'  {i+1:02d}   {item}';r.font.size=Pt(21);r.font.name='Inter';r.font.color.rgb=WHITE
    # number in gold
    r2=p.add_run()
txt(sl,'[Insert Bifacial PV Diagram / Photo]',7.5,1.5,5,4.5,sz=14,c=LGRAY,al=PP_ALIGN.CENTER)

# ===== S3: INTRODUCTION =====
sl=prs.slides.add_slide(BL);bg(sl)
ttl(sl,'Introduction & Motivation')
bullets(sl,[
    'Solar PV: fastest-growing electricity source globally \u2014 1,500+ GW installed by 2024 (IEA)',
    'India targets 500 GW non-fossil capacity by 2030; receives ~5,000 TWh solar energy/year',
    'Bifacial PV modules capture irradiance on both surfaces \u2192 5\u201330% higher yield vs monofacial',
    'First commercial bifacial panels: 1984 \u2014 projected to capture ~35% of global PV market by 2027',
    'Rear-side performance depends critically on tilt angle, ground albedo, and mounting height',
    'Improper parameter selection negates bifacial advantages \u2192 need for optimization frameworks',
    'Gap: No accessible, location-specific tool exists for combined parametric optimization'
],sz=20)

# ===== S4: LITERATURE =====
sl=prs.slides.add_slide(BL);bg(sl)
ttl(sl,'Literature Survey')
bullets(sl,[
    ('Guerrero-Lemus et al. [7]:','Comprehensive bifacial PV technology review; standardization needs identified'),
    ('Yusufoglu et al. [2]:','Annual energy yield analysis; up to 30% gain at 2m height; log-height dependence'),
    ('Dincer & Ozer [3]:','Parametric analysis of albedo, tilt, height, mounting; 21.2% rear gain for aluminum'),
    ('Sun et al. [10]:','Global perspective; combined tilt-albedo-latitude optimization across locations'),
    ('Ganesan et al. [4]:','n-type PERT bifacial under 5 albedo surfaces; 21.4% avg gain with aluminum'),
    ('Pelaez et al. [5]:','Comparison of 5 irradiance models; ~20% bifacial gain validated against field data'),
    ('Basak et al. [11]:','Tilt optimization specific to bifacial; optimal tilt differs from monofacial'),
],sz=18)

# ===== S5: GAPS =====
sl=prs.slides.add_slide(BL);bg(sl)
ttl(sl,'Research Gaps Identified')
bullets(sl,[
    ('1. Limited Combined Analysis \u2014','Studies evaluate tilt, albedo, height independently; combined interaction effects largely unexplored'),
    ('2. Location-Specific Gap \u2014','Most studies limited to European/American climates; India-specific validation scarce'),
    ('3. Accessibility Barrier \u2014','Optimization tools require PVsyst, SAM, or MATLAB licenses; inaccessible to many designers'),
    ('4. Single Optimization Objective \u2014','Few frameworks consider dual objectives: max total energy vs. max rear-side gain'),
    ('5. No Integrated Platform \u2014','No web-based tool combines real-time API data retrieval with parametric sweep analysis'),
],sz=20)

# ===== S6: PROBLEM STATEMENT =====
sl=prs.slides.add_slide(BL);bg(sl)
ttl(sl,'Problem Statement & Objectives')
txt(sl,'\u201cHow can the rear-side energy contribution and overall performance of bifacial PV systems be maximized by optimally selecting tilt angle, ground albedo, and mounting height for a given geographical location?\u201d',0.8,1.3,11.5,1.5,sz=22,c=WHITE,b=True,al=PP_ALIGN.CENTER,fn='Inter')
txt(sl,'Research Objectives',0.7,3.2,5,0.5,sz=22,c=GOLD,b=True)
bullets(sl,[
    'Develop mathematical model for bifacial PV (front + rear irradiance with view factors)',
    'Implement in MATLAB/Simulink for I-V and P-V characteristic simulation',
    'Build web-based frontend platform with NASA POWER API integration',
    'Perform parametric sweep to identify optimal configurations',
    'Validate framework through case study and literature comparison',
],l=0.7,tp=3.8,sz=18)

# ===== S7: MATH FRONT =====
sl=prs.slides.add_slide(BL);bg(sl)
ttl(sl,'Mathematical Formulation \u2014 Front Irradiance',sz=28)
eqs=[('Eq. (1)','G_tot = G_front + G_rear','Total effective irradiance on bifacial module'),
     ('Eq. (2)','beam_tilted = DNI \u00d7 max(0, cos \u03b8\u1d62)','Direct beam on tilted surface'),
     ('Eq. (3)','sky_vf = (1 + cos \u03b2) / 2','Sky view factor (Liu-Jordan model)'),
     ('Eq. (4)','gnd_vf = (1 \u2212 cos \u03b2) / 2','Ground view factor'),
     ('Eq. (5)','G_front = beam + DHI\u00b7sky_vf + GHI\u00b7\u03c1\u00b7gnd_vf','Front irradiance: beam + diffuse + reflected')]
tx=sl.shapes.add_textbox(Inches(0.7),Inches(1.3),Inches(11.5),Inches(5.8));tf=tx.text_frame;tf.word_wrap=True
for i,(lbl,eq,desc) in enumerate(eqs):
    p=tf.paragraphs[0] if i==0 else tf.add_paragraph();p.space_after=Pt(14)
    r1=p.add_run();r1.text=f'{lbl}:  ';r1.font.size=Pt(18);r1.font.bold=True;r1.font.name='Inter';r1.font.color.rgb=GOLD
    r2=p.add_run();r2.text=eq;r2.font.size=Pt(22);r2.font.name='Cambria Math';r2.font.color.rgb=WHITE;r2.font.italic=True
    r3=p.add_run();r3.text=f'\n         {desc}';r3.font.size=Pt(14);r3.font.name='Inter';r3.font.color.rgb=LGRAY

# ===== S8: MATH REAR =====
sl=prs.slides.add_slide(BL);bg(sl)
ttl(sl,'Mathematical Formulation \u2014 Rear Irradiance & Shadow',sz=26)
eqs=[('Eq. (6)','\u03b3_prof = arctan(tan \u03b1_sun / cos(\u03b3_sun \u2212 \u03b3_panel))','Solar profile angle'),
     ('Eq. (7)','shadow_lower = h / tan(\u03b3_prof)','Lower shadow boundary'),
     ('Eq. (8)','shadow_upper = (h+W\u00b7sin\u03b2) / tan(\u03b3_prof)','Upper shadow boundary'),
     ('Eq. (9)','F_V = \u222b cos\u03b8\u2081\u00b7cos\u03b8\u2082 / (2r) dx','Shadow view factor (Simpson\u2019s rule)'),
     ('Eq. (10)','G_r0 = \u03c1\u00b7DHI\u00b7rear_vf + \u03c1\u00b7(GHI\u2212DHI)\u00b7(rear_vf\u2212F_V)','Rear irradiance with shadow correction'),
     ('Eq. (11)','G_rear = G_r0 \u00d7 \u03c6_bifaciality','Scaled by bifaciality factor (0.7)')]
tx=sl.shapes.add_textbox(Inches(0.7),Inches(1.3),Inches(11.5),Inches(5.8));tf=tx.text_frame;tf.word_wrap=True
for i,(lbl,eq,desc) in enumerate(eqs):
    p=tf.paragraphs[0] if i==0 else tf.add_paragraph();p.space_after=Pt(11)
    r1=p.add_run();r1.text=f'{lbl}:  ';r1.font.size=Pt(17);r1.font.bold=True;r1.font.name='Inter';r1.font.color.rgb=GOLD
    r2=p.add_run();r2.text=eq;r2.font.size=Pt(20);r2.font.name='Cambria Math';r2.font.color.rgb=WHITE;r2.font.italic=True
    r3=p.add_run();r3.text=f'\n         {desc}';r3.font.size=Pt(13);r3.font.name='Inter';r3.font.color.rgb=LGRAY

# ===== S9: SIMULINK =====
sl=prs.slides.add_slide(BL);bg(sl)
ttl(sl,'Simulink Model & Implementation',sz=30)
txt(sl,'[Insert Simulink Model Screenshot]',0.5,1.5,6,4,sz=15,c=LGRAY,al=PP_ALIGN.CENTER)
bullets(sl,[
    'calculate_irradiance MATLAB function\nimplements Equations (1)\u2013(11)',
    'PV Array block: single-diode model\n550 Wp bifacial module, 72\u00d72 half-cells',
    'V_ramp sweep: 0V \u2192 Voc for complete\nI-V / P-V characteristic curves',
    'NOCT thermal model for cell\ntemperature derating',
    'Simpson\u2019s rule numerical integration\nfor shadow view factor F_V (N=200)',
],l=7,tp=1.5,w=5.5,sz=16)

# ===== S10: PLATFORM =====
sl=prs.slides.add_slide(BL);bg(sl)
ttl(sl,'Frontend Platform Architecture',sz=30)
layers=[('\u25b6 User Interface:','City, date range, sweep parameters, albedo surface type, optimization objective'),
        ('\u25b6 Backend (Node.js):','Express server, API routing, input validation, parametric sweep coordination'),
        ('\u25b6 NASA POWER API:','Hourly GHI, DHI, DNI, temperature \u2014 any global location, 0.5\u00b0 resolution'),
        ('\u25b6 Analysis Engine:','Erbs decomposition, Liu-Jordan transposition, view factor computation, I-V/P-V curves'),
        ('\u25b6 Results Dashboard:','Optimal config tables, charts, individual parameter variation analysis')]
tx=sl.shapes.add_textbox(Inches(0.7),Inches(1.3),Inches(11.5),Inches(5.5));tf=tx.text_frame;tf.word_wrap=True
for i,(t,d) in enumerate(layers):
    p=tf.paragraphs[0] if i==0 else tf.add_paragraph();p.space_after=Pt(14)
    r1=p.add_run();r1.text=t;r1.font.size=Pt(20);r1.font.bold=True;r1.font.name='Inter';r1.font.color.rgb=GOLD
    r2=p.add_run();r2.text=' '+d;r2.font.size=Pt(18);r2.font.name='Inter';r2.font.color.rgb=WHITE

prs.save('_endsem_fresh_p1.pptx')
print('Part 1 done (10 slides)')
