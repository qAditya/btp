import sys,io
sys.stdout=io.TextIOWrapper(sys.stdout.buffer,encoding='utf-8')
from pptx import Presentation
prs=Presentation('MidSem_PPT_BTP-2.pptx')
print(f'Total slides: {len(prs.slides)}')
print(f'Size: {prs.slide_width/914400:.1f} x {prs.slide_height/914400:.1f} inches')
for i,sl in enumerate(prs.slides):
    titles=[]
    for sh in sl.shapes:
        if hasattr(sh,'text') and sh.text.strip():
            t=sh.text.strip()[:80]
            if len(t)>3: titles.append(t)
    title = titles[0] if titles else 'no text'
    print(f'Slide {i+1}: {title}')
